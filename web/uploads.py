from __future__ import annotations

import base64
import csv
import io
import json
import subprocess
from dataclasses import dataclass
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import BadZipFile

import pypdfium2 as pdfium
import pytesseract
from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PresentationPackageNotFoundError
from pypdf import PdfReader
from pypdf.errors import PdfReadError
from PIL import Image, ImageOps, UnidentifiedImageError


MAX_DOCUMENT_UPLOAD_BYTES = 10 * 1024 * 1024
MAX_IMAGE_UPLOAD_BYTES = 20 * 1024 * 1024
MAX_VIDEO_UPLOAD_BYTES = 100 * 1024 * 1024
MAX_EXTRACTED_CHARS = 40_000
IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"})
VIDEO_EXTENSIONS = frozenset({".mp4", ".mov", ".webm", ".mkv", ".avi"})
SOURCE_EXTENSIONS = frozenset({
    ".txt", ".md", ".py", ".js", ".jsx", ".ts", ".tsx", ".html", ".css", ".xml",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".sh", ".sql", ".go", ".rs", ".java",
    ".kt", ".c", ".h", ".cpp", ".hpp",
})
SUPPORTED_EXTENSIONS = SOURCE_EXTENSIONS | frozenset({".csv", ".json", ".pdf", ".docx", ".xlsx", ".pptx"}) | IMAGE_EXTENSIONS | VIDEO_EXTENSIONS
MAX_MEDIA_IMAGES = 6
Image.MAX_IMAGE_PIXELS = 40_000_000


class UploadError(ValueError):
    pass


@dataclass(frozen=True)
class ExtractedUpload:
    text: str
    truncated: bool
    images: tuple[tuple[str, str, bytes], ...] = ()


def safe_filename(filename: str) -> str:
    return "".join(character for character in Path(filename).name if character.isprintable())[:160]


def max_upload_bytes(filename: str) -> int:
    suffix = Path(filename).suffix.lower()
    if suffix in VIDEO_EXTENSIONS:
        return MAX_VIDEO_UPLOAD_BYTES
    if suffix in IMAGE_EXTENSIONS:
        return MAX_IMAGE_UPLOAD_BYTES
    return MAX_DOCUMENT_UPLOAD_BYTES


def image_thumbnail_data_url(content: bytes) -> str:
    with Image.open(io.BytesIO(content)) as source:
        image = ImageOps.exif_transpose(source).convert("RGB")
        image.thumbnail((320, 320))
        stream = io.BytesIO()
        image.save(stream, format="JPEG", quality=80, optimize=True)
    return f"data:image/jpeg;base64,{base64.b64encode(stream.getvalue()).decode()}"


def extract_text(filename: str, content: bytes) -> ExtractedUpload:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise UploadError("unsupported file type; use text, document, spreadsheet, or image files")
    if not content:
        raise UploadError("file is empty")
    limit = max_upload_bytes(filename)
    if len(content) > limit:
        raise UploadError(f"file exceeds the {limit // (1024 * 1024)} MB limit")
    try:
        images: tuple[tuple[str, str, bytes], ...] = ()
        if suffix in SOURCE_EXTENSIONS:
            text = content.decode("utf-8-sig")
        elif suffix in IMAGE_EXTENSIONS:
            text, normalized_image = _extract_image(content)
            images = ((f"Uploaded image: {filename}", "image/jpeg", normalized_image),)
        elif suffix == ".csv":
            text = _extract_csv(content)
        elif suffix == ".json":
            text = json.dumps(json.loads(content.decode("utf-8-sig")), ensure_ascii=False, indent=2)
        elif suffix == ".pdf":
            text, images = _extract_pdf(content)
        elif suffix == ".docx":
            document = Document(io.BytesIO(content))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        elif suffix == ".xlsx":
            text = _extract_xlsx(content)
        elif suffix == ".pptx":
            text = _extract_pptx(content)
        else:
            text, images = _extract_video(filename, content)
    except UploadError:
        raise
    except (
        UnicodeDecodeError,
        csv.Error,
        json.JSONDecodeError,
        OSError,
        ValueError,
        KeyError,
        BadZipFile,
        PackageNotFoundError,
        PresentationPackageNotFoundError,
        InvalidFileException,
        PdfReadError,
        RuntimeError,
        subprocess.CalledProcessError,
        subprocess.TimeoutExpired,
        UnidentifiedImageError,
        Image.DecompressionBombError,
    ) as error:
        raise UploadError(f"could not read {suffix[1:].upper()} content") from error
    text = text.replace("\x00", "").strip()
    if not text:
        raise UploadError("file contains no extractable text")
    truncated = len(text) > MAX_EXTRACTED_CHARS
    return ExtractedUpload(text[:MAX_EXTRACTED_CHARS], truncated, images)


def _extract_image(content: bytes) -> tuple[str, bytes]:
    with Image.open(io.BytesIO(content)) as source:
        return _process_image(ImageOps.exif_transpose(source), "Image attached for visual analysis.")


def _process_image(source: Image.Image, description: str) -> tuple[str, bytes]:
    image = source.convert("RGB")
    image.thumbnail((2_048, 2_048))
    ocr_text = pytesseract.image_to_string(image, lang="kor+eng", timeout=30).strip()
    stream = io.BytesIO()
    image.save(stream, format="JPEG", quality=88, optimize=True)
    text = description
    if ocr_text:
        text += f"\nOCR text:\n{ocr_text}"
    else:
        text += " OCR found no text."
    return text, stream.getvalue()


def _extract_pdf(content: bytes) -> tuple[str, tuple[tuple[str, str, bytes], ...]]:
    reader = PdfReader(io.BytesIO(content))
    page_texts = [(page.extract_text() or "").strip() for page in reader.pages]
    scanned_pages = [index for index, text in enumerate(page_texts) if len(text) < 20][:MAX_MEDIA_IMAGES]
    images: list[tuple[str, str, bytes]] = []
    if scanned_pages:
        document = pdfium.PdfDocument(content)
        try:
            for page_index in scanned_pages:
                rendered = document[page_index].render(scale=2).to_pil()
                text, normalized = _process_image(rendered, f"Scanned PDF page {page_index + 1}.")
                page_texts[page_index] = text
                images.append((f"Scanned PDF page {page_index + 1}", "image/jpeg", normalized))
        finally:
            document.close()
    pages = "\n\n".join(f"Page {index + 1}:\n{text}" for index, text in enumerate(page_texts))
    return pages, tuple(images)


def _extract_video(filename: str, content: bytes) -> tuple[str, tuple[tuple[str, str, bytes], ...]]:
    suffix = Path(filename).suffix.lower()
    with TemporaryDirectory(prefix="local-ai-video-") as directory:
        video_path = Path(directory) / f"input{suffix}"
        video_path.write_bytes(content)
        probe = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=nw=1:nk=1", str(video_path)],
            check=True,
            capture_output=True,
            text=True,
            timeout=20,
        )
        duration = float(probe.stdout.strip())
        if duration <= 0 or duration > 7_200:
            raise UploadError("video duration must be between 1 second and 2 hours")
        interval = max(duration / MAX_MEDIA_IMAGES, 1.0)
        output_pattern = str(Path(directory) / "frame-%02d.jpg")
        subprocess.run(
            [
                "ffmpeg", "-v", "error", "-i", str(video_path),
                "-vf", f"fps=1/{interval},scale='min(1600,iw)':-2", "-frames:v", str(MAX_MEDIA_IMAGES),
                "-q:v", "3", output_pattern,
            ],
            check=True,
            capture_output=True,
            timeout=60,
        )
        frames = sorted(Path(directory).glob("frame-*.jpg"))
        if not frames:
            raise UploadError("video contains no readable frames")
        texts = [f"Video duration: {duration:.1f} seconds. Representative frames follow."]
        images: list[tuple[str, str, bytes]] = []
        for index, frame_path in enumerate(frames):
            timestamp = min(index * interval, duration)
            with Image.open(frame_path) as frame:
                text, normalized = _process_image(frame, f"Video frame near {timestamp:.1f} seconds.")
            texts.append(text)
            images.append((f"Video frame near {timestamp:.1f} seconds", "image/jpeg", normalized))
    return "\n\n".join(texts), tuple(images)


def _extract_csv(content: bytes) -> str:
    stream = io.StringIO(content.decode("utf-8-sig"))
    rows: list[str] = []
    for row in csv.reader(stream):
        rows.append(" | ".join(cell.strip() for cell in row))
        if sum(len(item) for item in rows) >= MAX_EXTRACTED_CHARS:
            break
    return "\n".join(rows)


def _extract_xlsx(content: bytes) -> str:
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines: list[str] = []
    try:
        for sheet in workbook.worksheets:
            lines.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) if value is not None else "" for value in row]
                lines.append(" | ".join(values))
                if sum(len(item) for item in lines) >= MAX_EXTRACTED_CHARS:
                    return "\n".join(lines)
    finally:
        workbook.close()
    return "\n".join(lines)


def _extract_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        slides.append(f"Slide {index}:\n{text}")
    return "\n\n".join(slides)